import io
import json
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def sanitize_name(name: str) -> str:
    return "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in name)


def extract_images_from_pdf(pdf_path: Path, course: str, output_dir: Path) -> list:
    """
    Extract images from a PDF file and save them to output_dir.
    Returns a list of metadata dictionaries.
    """
    metadata = []
    doc = fitz.open(pdf_path)

    base_name = sanitize_name(pdf_path.stem)

    for page_index in range(len(doc)):
        page = doc[page_index]
        image_list = page.get_images(full=True)

        for img_num, img in enumerate(image_list, start=1):
            xref = img[0]
            image_info = doc.extract_image(xref)
            image_bytes = image_info["image"]
            image_ext = image_info["ext"]

            img = Image.open(io.BytesIO(image_bytes))
            width, height = img.size

            if width < 100 or height < 100:
              continue
            image_filename = f"{base_name}_page{page_index + 1}_img{img_num}.{image_ext}"
            image_path = output_dir / image_filename

            with open(image_path, "wb") as f:
                f.write(image_bytes)

            metadata.append({
                "course": course,
                "source_file": pdf_path.name,
                "file_type": "pdf",
                "page_or_slide": page_index + 1,
                "image_index": img_num,
                "image_path": str(image_path)
            })

    doc.close()
    return metadata


def extract_images_from_pptx(pptx_path: Path, course: str, output_dir: Path) -> list:
    """
    Extract images from a PPTX file and save them to output_dir.
    Returns a list of metadata dictionaries.
    """
    metadata = []
    prs = Presentation(pptx_path)

    base_name = sanitize_name(pptx_path.stem)

    for slide_index, slide in enumerate(prs.slides, start=1):
        img_num = 0

        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img_num += 1
                image = shape.image
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                width , height = img.size

                if width < 100 or height < 100:
                 continue
                image_ext = image.ext

                image_filename = f"{base_name}_slide{slide_index}_img{img_num}.{image_ext}"
                image_path = output_dir / image_filename

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                metadata.append({
                    "course": course,
                    "source_file": pptx_path.name,
                    "file_type": "pptx",
                    "page_or_slide": slide_index,
                    "image_index": img_num,
                    "image_path": str(image_path)
                })

    return metadata


def extract_images_for_course(course_dir: Path, processed_dir: Path) -> list:
    """
    Extract images from all PDF/PPTX files inside one course folder.
    Saves images to data/processed/<course>/images/
    Saves metadata to data/processed/<course>/image_metadata.json
    """
    course = course_dir.name
    images_output_dir = processed_dir / course / "images"
    ensure_dir(images_output_dir)

    all_metadata = []

    for file_path in course_dir.iterdir():
        if file_path.is_file():
            suffix = file_path.suffix.lower()

            if suffix == ".pdf":
                print(f"Extracting images from PDF: {file_path.name}")
                all_metadata.extend(
                    extract_images_from_pdf(file_path, course, images_output_dir)
                )

            elif suffix == ".pptx":
                print(f"Extracting images from PPTX: {file_path.name}")
                all_metadata.extend(
                    extract_images_from_pptx(file_path, course, images_output_dir)
                )

    metadata_path = processed_dir / course / "image_metadata.json"
    ensure_dir(metadata_path.parent)

    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(all_metadata, f, indent=2, ensure_ascii=False)

    print(f"Saved metadata to: {metadata_path}")
    print(f"Extracted {len(all_metadata)} images for course {course}")

    return all_metadata


def main():
    project_root = Path(__file__).resolve().parent.parent
    raw_data_dir = project_root / "data" / "raw"
    processed_data_dir = project_root / "data" / "processed"

    if not raw_data_dir.exists():
        print(f"Raw data directory not found: {raw_data_dir}")
        return

    for course_dir in raw_data_dir.iterdir():
        if course_dir.is_dir():
            print(f"\nProcessing course folder: {course_dir.name}")
            extract_images_for_course(course_dir, processed_data_dir)


if __name__ == "__main__":
    main()