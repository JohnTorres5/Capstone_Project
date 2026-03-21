import os
import json
import re
from pathlib import Path
import fitz 
from pptx import Presentation

# Text Cleaning 
def clean_text(text: str) -> str:
    text = re.sub(r'\s+', ' ', text)
    text = text.replace('\x00', '')
    return text.strip()

# PDF Extraction
def extract_text_from_pdf(file_path: str):
    doc = fitz.open(file_path)
    full_text = ""
    pages_list = []

    for i, page in enumerate(doc, start=1):
        page_text = page.get_text()
        cleaned = clean_text(page_text)

        pages_list.append({
            "page": i,
            "text": cleaned
        })

        full_text += cleaned + "\n"

    doc.close()
    return full_text.strip(), pages_list

# PPTX Extraction
def extract_text_from_pptx(file_path: str):
    presentation = Presentation(file_path)
    full_text = ""
    pages_list = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

        cleaned = clean_text(slide_text)

        pages_list.append({
            "slide": slide_number,
            "text": cleaned
        })

        full_text += f"\n--- Slide {slide_number} ---\n{cleaned}"

    return full_text.strip(), pages_list

# Preprocessing Function
def preprocess_course(input_dir: str, output_dir: str, course_name: str):
    os.makedirs(output_dir, exist_ok=True)
    input_path = Path(input_dir)

    for file_path in input_path.glob("*"):
        if file_path.suffix.lower() not in [".pdf", ".pptx"]:
            continue

        print(f"Processing: {file_path.name}")

        # Extract text
        if file_path.suffix.lower() == ".pdf":
            raw_text, pages_list = extract_text_from_pdf(str(file_path))
        elif file_path.suffix.lower() == ".pptx":
            raw_text, pages_list = extract_text_from_pptx(str(file_path))

        # Clean full text (optional, already cleaned per page)
        cleaned_text = clean_text(raw_text)

        # Create structured output
        document_data = {
            "course": course_name,
            "source_file": file_path.name,
            "text": cleaned_text,
            "pages": pages_list
        }

        # Save as JSON
        output_file = Path(output_dir) / f"{file_path.stem}.json"
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(document_data, f, indent=2)

        print(f"Saved: {output_file.name}")

    print(f"\nFinished preprocessing for course: {course_name}")